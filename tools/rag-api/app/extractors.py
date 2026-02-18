"""Text extraction utilities for different document types."""

import logging
from io import BytesIO
from typing import List, Tuple, Optional
import re

try:
    from pypdf import PdfReader
except ImportError:
    from PyPDF2 import PdfReader

try:
    import docx
except ImportError:
    docx = None

try:
    import markdown
except ImportError:
    markdown = None

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

logger = logging.getLogger(__name__)


def clean_text(text: str) -> str:
    """
    Clean and normalize extracted text.
    
    Args:
        text: Raw extracted text
    
    Returns:
        Cleaned text string
    """
    if not text:
        return ""
    
    # Remove null characters
    text = text.replace("\x00", " ")
    
    # Normalize whitespace
    text = re.sub(r'\s+', ' ', text)
    
    # Remove excessive newlines but preserve paragraph breaks
    text = re.sub(r'\n\s*\n\s*\n+', '\n\n', text)
    
    # Clean up common OCR artifacts
    text = re.sub(r'[^\w\s\-.,;:!?()[\]{}"\'/\\@#$%^&*+=<>~`|]', ' ', text)
    
    return text.strip()


def extract_pdf_text(binary_content: bytes) -> List[Tuple[int, str]]:
    """
    Extract text from PDF file.
    
    Args:
        binary_content: PDF file content as bytes
    
    Returns:
        List of tuples containing (page_number, text_content)
    """
    pages = []
    
    try:
        reader = PdfReader(BytesIO(binary_content))
        
        for page_num, page in enumerate(reader.pages, start=1):
            try:
                text = page.extract_text() or ""
                cleaned_text = clean_text(text)
                
                if cleaned_text:  # Only add pages with actual content
                    pages.append((page_num, cleaned_text))
                    
            except Exception as e:
                logger.warning(f"Failed to extract text from page {page_num}: {e}")
                continue
                
    except Exception as e:
        logger.error(f"Failed to parse PDF: {e}")
        raise ValueError(f"Unable to parse PDF: {e}")
    
    return pages


def extract_docx_text(binary_content: bytes) -> str:
    """
    Extract text from DOCX file.
    
    Args:
        binary_content: DOCX file content as bytes
    
    Returns:
        Extracted text content
    """
    if docx is None:
        raise ImportError("python-docx package is required for DOCX extraction")
    
    try:
        document = docx.Document(BytesIO(binary_content))
        
        # Extract text from paragraphs
        paragraphs = []
        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if text:
                paragraphs.append(text)
        
        # Extract text from tables
        for table in document.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text:
                        row_text.append(cell_text)
                if row_text:
                    paragraphs.append(" | ".join(row_text))
        
        full_text = "\n".join(paragraphs)
        return clean_text(full_text)
        
    except Exception as e:
        logger.error(f"Failed to parse DOCX: {e}")
        raise ValueError(f"Unable to parse DOCX: {e}")


def extract_txt_text(binary_content: bytes) -> str:
    """
    Extract text from plain text file.
    
    Args:
        binary_content: Text file content as bytes
    
    Returns:
        Extracted text content
    """
    try:
        # Try different encodings
        encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
        
        for encoding in encodings:
            try:
                text = binary_content.decode(encoding)
                return clean_text(text)
            except UnicodeDecodeError:
                continue
        
        # If all encodings fail, use utf-8 with error handling
        text = binary_content.decode('utf-8', errors='replace')
        return clean_text(text)
        
    except Exception as e:
        logger.error(f"Failed to decode text file: {e}")
        raise ValueError(f"Unable to decode text file: {e}")


def extract_markdown_text(binary_content: bytes) -> str:
    """
    Extract text from Markdown file.
    
    Args:
        binary_content: Markdown file content as bytes
    
    Returns:
        Extracted text content (markdown converted to plain text)
    """
    try:
        # First decode as text
        text = extract_txt_text(binary_content)
        
        # If markdown library is available, convert to HTML then strip tags
        if markdown and BeautifulSoup:
            html = markdown.markdown(text)
            soup = BeautifulSoup(html, 'html.parser')
            plain_text = soup.get_text(separator=' ')
            return clean_text(plain_text)
        else:
            # Fallback: just return the raw markdown as text
            # Strip common markdown syntax for better readability
            text = re.sub(r'#{1,6}\s+', '', text)  # Remove headers
            text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)  # Remove bold
            text = re.sub(r'\*(.+?)\*', r'\1', text)  # Remove italic
            text = re.sub(r'`(.+?)`', r'\1', text)  # Remove code
            text = re.sub(r'\[(.+?)\]\(.+?\)', r'\1', text)  # Remove links, keep text
            return clean_text(text)
        
    except Exception as e:
        logger.error(f"Failed to parse Markdown: {e}")
        raise ValueError(f"Unable to parse Markdown: {e}")


def extract_html_text(binary_content: bytes) -> str:
    """
    Extract text from HTML file.
    
    Args:
        binary_content: HTML file content as bytes
    
    Returns:
        Extracted text content
    """
    if BeautifulSoup is None:
        raise ImportError("beautifulsoup4 package is required for HTML extraction")
    
    try:
        # Decode HTML
        encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
        html_content = None
        
        for encoding in encodings:
            try:
                html_content = binary_content.decode(encoding)
                break
            except UnicodeDecodeError:
                continue
        
        if not html_content:
            html_content = binary_content.decode('utf-8', errors='replace')
        
        # Parse HTML and extract text
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()
        
        # Get text
        text = soup.get_text(separator=' ')
        return clean_text(text)
        
    except Exception as e:
        logger.error(f"Failed to parse HTML: {e}")
        raise ValueError(f"Unable to parse HTML: {e}")


def extract_xlsx_text(binary_content: bytes) -> str:
    """
    Extract text from Excel file.
    
    Args:
        binary_content: Excel file content as bytes
    
    Returns:
        Extracted text content
    """
    if openpyxl is None:
        raise ImportError("openpyxl package is required for Excel extraction")
    
    try:
        workbook = openpyxl.load_workbook(BytesIO(binary_content), data_only=True)
        
        all_text = []
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            all_text.append(f"Sheet: {sheet_name}")
            
            for row in sheet.iter_rows(values_only=True):
                row_values = [str(cell) if cell is not None else "" for cell in row]
                row_text = " | ".join(row_values).strip()
                if row_text:
                    all_text.append(row_text)
        
        return clean_text("\n".join(all_text))
        
    except Exception as e:
        logger.error(f"Failed to parse Excel: {e}")
        raise ValueError(f"Unable to parse Excel: {e}")


def extract_pptx_text(binary_content: bytes) -> str:
    """
    Extract text from PowerPoint file.
    
    Args:
        binary_content: PowerPoint file content as bytes
    
    Returns:
        Extracted text content
    """
    if Presentation is None:
        raise ImportError("python-pptx package is required for PowerPoint extraction")
    
    try:
        prs = Presentation(BytesIO(binary_content))
        
        all_text = []
        for i, slide in enumerate(prs.slides, start=1):
            all_text.append(f"Slide {i}:")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    all_text.append(shape.text)
        
        return clean_text("\n".join(all_text))
        
    except Exception as e:
        logger.error(f"Failed to parse PowerPoint: {e}")
        raise ValueError(f"Unable to parse PowerPoint: {e}")


def detect_file_type(filename: str, content: bytes) -> str:
    """
    Detect file type based on filename and content.
    
    Args:
        filename: Original filename
        content: File content as bytes
    
    Returns:
        File type string ('pdf', 'docx', 'txt', 'md', 'html', 'xlsx', 'pptx', 'unknown')
    """
    filename_lower = filename.lower()
    
    # Check by file extension first
    if filename_lower.endswith('.pdf'):
        return 'pdf'
    elif filename_lower.endswith('.docx'):
        return 'docx'
    elif filename_lower.endswith(('.txt', '.text')):
        return 'txt'
    elif filename_lower.endswith(('.md', '.markdown')):
        return 'md'
    elif filename_lower.endswith(('.html', '.htm')):
        return 'html'
    elif filename_lower.endswith(('.xlsx', '.xlsm')):
        return 'xlsx'
    elif filename_lower.endswith('.xls'):
        return 'xls'  # Old Excel format
    elif filename_lower.endswith('.pptx'):
        return 'pptx'
    elif filename_lower.endswith('.ppt'):
        return 'ppt'  # Old PowerPoint format
    
    # Check by file signature (magic bytes)
    if content.startswith(b'%PDF'):
        return 'pdf'
    elif content.startswith(b'PK\x03\x04'):
        # This is a ZIP-based format, check what's inside
        if b'word/' in content[:2048]:
            return 'docx'
        elif b'xl/' in content[:2048]:
            return 'xlsx'
        elif b'ppt/' in content[:2048]:
            return 'pptx'
    
    # Check for HTML
    if b'<html' in content[:1024].lower() or b'<!doctype html' in content[:1024].lower():
        return 'html'
    
    # Default to txt for other text-like content
    try:
        content[:1024].decode('utf-8')
        return 'txt'
    except UnicodeDecodeError:
        pass
    
    return 'unknown'


def extract_text_from_file(filename: str, binary_content: bytes) -> List[Tuple[Optional[int], str]]:
    """
    Extract text from a file based on its type.
    
    Args:
        filename: Original filename
        binary_content: File content as bytes
    
    Returns:
        List of tuples containing (page_number, text_content)
        For non-paginated formats, page_number will be None
    """
    file_type = detect_file_type(filename, binary_content)
    
    try:
        if file_type == 'pdf':
            return extract_pdf_text(binary_content)
        elif file_type == 'docx':
            text = extract_docx_text(binary_content)
            return [(None, text)] if text else []
        elif file_type == 'txt':
            text = extract_txt_text(binary_content)
            return [(None, text)] if text else []
        elif file_type == 'md':
            text = extract_markdown_text(binary_content)
            return [(None, text)] if text else []
        elif file_type == 'html':
            text = extract_html_text(binary_content)
            return [(None, text)] if text else []
        elif file_type in ('xlsx', 'xls'):
            text = extract_xlsx_text(binary_content)
            return [(None, text)] if text else []
        elif file_type in ('pptx', 'ppt'):
            text = extract_pptx_text(binary_content)
            return [(None, text)] if text else []
        else:
            logger.warning(f"Unsupported file type '{file_type}' for file '{filename}'")
            return []
            
    except ImportError as e:
        logger.warning(f"Missing library for {file_type}: {e}. File '{filename}' will be skipped.")
        return []
    except Exception as e:
        logger.error(f"Failed to extract text from {filename}: {e}")
        return []


def get_supported_extensions() -> List[str]:
    """
    Get list of supported file extensions.
    
    Returns:
        List of supported file extensions
    """
    extensions = ['.pdf', '.txt', '.text', '.md', '.markdown', '.html', '.htm']
    
    # Add extensions for optional dependencies
    if docx:
        extensions.append('.docx')
    if openpyxl:
        extensions.extend(['.xlsx', '.xlsm', '.xls'])
    if Presentation:
        extensions.extend(['.pptx', '.ppt'])
    
    return extensions
